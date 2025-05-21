import os
import json
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from bs4 import BeautifulSoup
from ebooklib import epub

def extract_text_from_pdf(pfad):
    texte = []
    try:
        doc = fitz.open(pfad)
        for seite in doc:
            text = seite.get_text()
            if text:
                texte.append(text)
    except Exception as e:
        raise RuntimeError(f"PDF-Leseproblem: {e}")
    return "\n".join(texte)
#    # Diese Funktion extrahiert den Text aus einer PDF-Datei

def extract_text_from_docx(pfad):
    try:
        doc = docx.Document(pfad)
        return "\n".join(para.text for para in doc.paragraphs if para.text.strip())
    except Exception as e:
        raise RuntimeError(f"DOCX-Leseproblem: {e}")
#    # Diese Funktion extrahiert den Text aus einer DOCX-Datei

def extract_text_from_pptx(pfad):
    try:
        prs = Presentation(pfad)
        texte = []
        for folie in prs.slides:
            for shape in folie.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        texte.append(text)
        return "\n".join(texte)
    except Exception as e:
        raise RuntimeError(f"PPTX-Leseproblem: {e}")
#    # Diese Funktion extrahiert den Text aus einer PPTX-Datei

def extract_text_from_html(pfad):
    try:
        with open(pfad, encoding="utf-8") as f:
            soup = BeautifulSoup(f, "lxml")
        return soup.get_text()
    except Exception as e:
        raise RuntimeError(f"HTML-Leseproblem: {e}")
#    # Diese Funktion extrahiert den Text aus einer HTML-Datei

def extract_text_from_epub(pfad):
    try:
        book = epub.read_epub(pfad)
        texte = []
        for item in book.get_items():
            if item.get_type() == epub.EpubHtml:
                soup = BeautifulSoup(item.get_content(), "lxml")
                texte.append(soup.get_text())
        return "\n".join(texte)
    except Exception as e:
        raise RuntimeError(f"EPUB-Leseproblem: {e}")
#    # Diese Funktion extrahiert den Text aus einer EPUB-Datei

def extract_text_from_txt(pfad):
    try:
        with open(pfad, encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        raise RuntimeError(f"TXT-Leseproblem: {e}")
#    # Diese Funktion extrahiert den Text aus einer TXT-Datei

def extract_text_from_rtf(pfad):
    try:
        import striprtf
        with open(pfad, encoding="utf-8") as f:
            rtf_content = f.read()
        return striprtf.rtf_to_text(rtf_content)
    except ImportError:
        raise RuntimeError("Das Paket 'striprtf' wird für RTF benötigt. Installiere es mit 'pip install striprtf'.")
    except Exception as e:
        raise RuntimeError(f"RTF-Leseproblem: {e}")
#    # Diese Funktion extrahiert den Text aus einer RTF-Datei

def lade_dokumente_aus_ordner(ordnerpfad, fach=None):
    texte = []
    fehler = []

    for wurzel, _, dateien in os.walk(ordnerpfad):
        for datei in dateien:
            pfad = os.path.join(wurzel, datei)
            try:
                if datei.endswith(".pdf"):
                    text = extract_text_from_pdf(pfad)
                elif datei.endswith(".docx"):
                    text = extract_text_from_docx(pfad)
                elif datei.endswith(".pptx"):
                    text = extract_text_from_pptx(pfad)
                elif datei.endswith((".html", ".htm")):
                    text = extract_text_from_html(pfad)
                elif datei.endswith(".epub"):
                    text = extract_text_from_epub(pfad)
                elif datei.endswith(".txt"):
                    text = extract_text_from_txt(pfad)
                elif datei.endswith(".rtf"):
                    text = extract_text_from_rtf(pfad)
                else:
                    continue
                texte.append({"text": text, "quelle": pfad})
            except Exception as e:
                fehler.append({"datei": pfad, "fehler": str(e)})

    # Fehler in JSON speichern, wenn Fach übergeben
    if fach:
        fehler_pfad = os.path.join("data", "cache", f"fehler_{fach}.json")
        os.makedirs(os.path.dirname(fehler_pfad), exist_ok=True)
        with open(fehler_pfad, "w", encoding="utf-8") as f:
            json.dump(fehler, f, indent=2, ensure_ascii=False)

    return texte
